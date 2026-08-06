"""Text extraction for supported course-material formats.

`extract_text` dispatches on MIME type first (falling back to the path's
file extension when the MIME type is missing or unrecognized -- blob-store
paths are sha256-named with no extension, so in the real pipeline MIME is
what actually drives dispatch; the extension fallback mainly serves callers
that hand in a real filename, including these tests).

Never raises: every per-format extractor is called inside a broad
try/except here, so malformed input (a truncated PDF, garbled bytes given a
misleading MIME type, etc.) becomes a `None` result rather than an
exception. Callers (the S1 summarize stage's extract pass) turn a `None`
into `status='failed', error='unsupported-or-unparseable'`.
"""

from __future__ import annotations

import re
from collections.abc import Callable
from pathlib import Path

_EXCESS_NEWLINES_RE = re.compile(r"\n{3,}")
_CUE_NUMBER_RE = re.compile(r"^\d+$")

# MIME type (main type/subtype, no parameters) -> internal format key.
_MIME_TO_FORMAT: dict[str, str] = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "text/html": "html",
    "text/vtt": "vtt",
    "application/x-subrip": "srt",
    "text/x-srt": "srt",
    "text/plain": "txt",
    "text/markdown": "md",
}

# File extension (lowercase, no dot) -> internal format key. Used only when
# the MIME type is missing or not in the table above.
_EXTENSION_TO_FORMAT: dict[str, str] = {
    "pdf": "pdf",
    "pptx": "pptx",
    "docx": "docx",
    "html": "html",
    "htm": "html",
    "vtt": "vtt",
    "srt": "srt",
    "txt": "txt",
    "md": "md",
    "markdown": "md",
}


def extract_text(path: Path, mime: str | None, kind: str) -> str | None:
    """Extract plain text from the file at `path`, or `None` if the format
    isn't supported or the content couldn't be parsed. `kind` (the
    material's guessed kind, e.g. 'slides'/'transcript') is accepted for
    parity with callers but dispatch is driven entirely by `mime`/extension,
    per format:

    - pdf: pymupdf, all pages' get_text() joined with form feeds
    - pptx: python-pptx, per-slide shape text + speaker notes, "[Slide N]"
    - docx: python-docx, paragraph text joined with newlines
    - html: BeautifulSoup get_text(" ", strip=True)
    - vtt/srt: strip the WEBVTT header, cue numbers, and "-->" timestamp
      lines; join remaining cue text, deduping consecutive identical lines
    - txt/md: read as UTF-8, replacing invalid bytes
    - anything else (video, images, unknown binaries): None
    """
    del kind  # not used for dispatch; see docstring

    fmt = _detect_format(path, mime)
    if fmt is None:
        return None

    extractor = _EXTRACTORS[fmt]
    try:
        text = extractor(path)
    except Exception:
        return None

    if text is None:
        return None
    return _normalize(text)


def _detect_format(path: Path, mime: str | None) -> str | None:
    if mime:
        bare_mime = mime.split(";", 1)[0].strip().lower()
        fmt = _MIME_TO_FORMAT.get(bare_mime)
        if fmt:
            return fmt

    ext = path.suffix.lstrip(".").lower()
    return _EXTENSION_TO_FORMAT.get(ext)


def _normalize(text: str) -> str:
    return _EXCESS_NEWLINES_RE.sub("\n\n", text)


# --------------------------------------------------------------------------
# Per-format extractors -- each returns the raw extracted text (normalization
# happens once, centrally, in extract_text) or raises on malformed input.
# --------------------------------------------------------------------------


def _extract_pdf(path: Path) -> str:
    import fitz  # PyMuPDF

    doc = fitz.open(path)
    try:
        return "\f".join(page.get_text() for page in doc)
    finally:
        doc.close()


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    blocks: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        shape_texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text
        ]
        block = f"[Slide {index}]\n" + "\n".join(shape_texts)

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes:
                block += f"\nNotes: {notes}"

        blocks.append(block)
    return "\n\n".join(blocks)


def _extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_html(path: Path) -> str:
    from bs4 import BeautifulSoup

    raw = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "html.parser")
    return soup.get_text(" ", strip=True)


def _extract_subtitle(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="replace")

    cue_lines: list[str] = []
    for raw_line in raw.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if line.upper().startswith("WEBVTT"):
            continue
        if "-->" in line:
            continue
        if _CUE_NUMBER_RE.match(line):
            continue
        if cue_lines and cue_lines[-1] == line:
            continue  # dedupe consecutive identical cue text
        cue_lines.append(line)

    return "\n".join(cue_lines)


def _extract_plain_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


_EXTRACTORS: dict[str, Callable[[Path], str]] = {
    "pdf": _extract_pdf,
    "pptx": _extract_pptx,
    "docx": _extract_docx,
    "html": _extract_html,
    "vtt": _extract_subtitle,
    "srt": _extract_subtitle,
    "txt": _extract_plain_text,
    "md": _extract_plain_text,
}
