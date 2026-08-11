"""A fake Brightspace (D2L Valence API) tenant for offline E2E testing.

Serves the subset of the real D2L REST surface the extension's sync engine
actually calls (see extension/src/lib/d2l-client.ts): API version discovery,
whoami, paged enrollments, a course's table of contents, per-topic file
download, and the news/dropbox "extras" the sync engine folds into /toc.

Lives in tests/ (not production code): importable by both pytest
(test_fake_d2l.py) and scripts/e2e.py. Nothing under backend/src ever
imports this module.

Two fault-injection modes, both driven off one shared per-app request
counter (every route except /_control/* ticks it):

- `rate_limit_429_every=N`: every Nth request gets a 429 (Retry-After: 1
  plus a couple of D2L-shaped rate-limit headers) instead of its real
  response.
- `expire_session_after=N`: once more than N requests have been made, EVERY
  route (except /_control/*) returns 401 -- simulating a dead Brightspace
  session cookie -- until `POST /_control/reset` is called.
"""

from __future__ import annotations

import io
import json
from functools import lru_cache
from pathlib import Path, PurePosixPath
from typing import Any, Callable
from urllib.parse import urlparse

from fastapi import FastAPI, HTTPException, Request, Response
from fastapi.responses import JSONResponse

# --------------------------------------------------------------------------
# Fixture tenant data
# --------------------------------------------------------------------------

ORG_UNIT_ID = 555
LP_VERSION = "1.35"
LE_VERSION = "1.79"

WHOAMI: dict[str, Any] = {"Identifier": "999", "UniqueName": "gavin", "DisplayName": "Gavin Fake"}

# Two enrollments on purpose: exercises both bookmark paging (page size 1,
# below) and the extension's course-type filter (isCourseEnrollment in
# d2l-client.ts matches Type.Code/Type.Name against /course/i -- the
# Department entry does not and must be dropped client-side).
ENROLLMENTS: list[dict[str, Any]] = [
    {
        "OrgUnit": {
            "Id": ORG_UNIT_ID,
            "Name": "Intro to CS",
            "Code": "CS101",
            "Type": {"Id": 3, "Code": "Course Offering", "Name": "Course Offering"},
        }
    },
    {
        "OrgUnit": {
            "Id": 9999,
            "Name": "Computer Science Department",
            "Code": None,
            "Type": {"Id": 1, "Code": "Department", "Name": "Department"},
        }
    },
]
_ENROLLMENT_PAGE_SIZE = 1

# News/dropbox items in the REAL Valence shape: PascalCase, with the body
# nested under `Body`/`CustomInstructions` as a {Text, Html} pair. They are
# deliberately NOT pre-shaped to the backend's /toc `extras` contract (see
# api/ingest.py's NewsExtra/DropboxExtra) -- the reshaping is the
# extension's job (d2l-client.ts's toNewsExtras/toDropboxExtras), and a
# fixture that hands out already-correct objects would make three layers of
# tests agree on a shape D2L never actually sends.
#
# The second news item exercises the defensive half of that mapping: a
# tenant that publishes an announcement with no HTML rendering at all.
ADMIN_NEWS_TITLE = "Office hours moved"
"""M3.5a: the fixture's administrative material. Its title carries one of
MockBackend's `_MOCK_ADMIN_TITLE_MARKERS` (agents/llm.py), so an offline run
classifies it `is_administrative=True` and S4 files it under the synthetic
"Logistics & admin" bucket -- which is what lets scripts/e2e.py assert that
bucket end to end instead of only unit tests setting the column by hand.
Named rather than inlined below because e2e.py asserts on it by title:
renaming the announcement without renaming this constant would quietly
delete that coverage."""

LTI_HINT_MEDIASITE_TITLE = "Mediasite Channel (Week 2 Lectures)"
"""M2.7: title of `toc_sample.json`'s TOPIC-1007, one of two Link topics
shaped exactly like an LTI-embedded recording channel (a `type=lti`
quicklink `Url` plus a title `_LTI_HINT_TITLE_RE`/api/media.py matches) --
the real-world shape this whole milestone started from. scripts/e2e.py
resolves this one to a mediasite-shaped `finalUrl` and asserts the
`resolved` round-trip. Named (not inlined) for the same reason as
`ADMIN_NEWS_TITLE` above: e2e.py looks it up by title in the
`lti-candidates` response, since the material's DB id isn't known ahead of
a sync."""

LTI_HINT_ZOOM_TITLE = "Zoom Recordings (Week 2 Lab)"
"""M2.7: title of `toc_sample.json`'s TOPIC-1008, the second LTI-shaped Link
topic -- see `LTI_HINT_MEDIASITE_TITLE` above. scripts/e2e.py resolves this
one to a URL `classify_url` does not recognize and asserts the
`unrecognized` round-trip."""

NEWS: list[dict[str, Any]] = [
    {
        "Id": 1,
        "Title": "Midterm date announced",
        "Body": {
            "Text": "The midterm will be held in week 6.",
            "Html": "<p>The midterm will be held in week 6.</p>",
        },
        "StartDate": "2026-01-05T12:00:00.000Z",
        "EndDate": None,
        "IsGlobal": False,
        "IsPublished": True,
    },
    {
        "Id": 2,
        "Title": ADMIN_NEWS_TITLE,
        "Body": {"Text": "Office hours are now on Thursdays.", "Html": None},
        "StartDate": "2026-01-06T12:00:00.000Z",
        "EndDate": None,
        "IsGlobal": False,
        "IsPublished": True,
    },
    # The third item exists for the E2E's media stage: an announcement whose
    # HTML body links a lecture recording is exactly how recordings show up
    # in a real course (they live outside D2L; the ToC only ever holds the
    # link), and it's what media/detect.py's HTML-page scan is written
    # against. The URL's "mock-captions" token is what makes
    # MockMediaFetcher return a caption track offline -- see its docstring.
    {
        "Id": 3,
        "Title": "Lecture 1 recording posted",
        "Body": {
            "Text": "The Week 1 lecture recording is up. Passcode: wk1pass",
            "Html": (
                '<p>The Week 1 lecture recording is up: '
                '<a href="https://zoom.us/rec/share/mock-captions-week1">watch it here</a>. '
                "Passcode: wk1pass</p>"
            ),
        },
        "StartDate": "2026-01-07T12:00:00.000Z",
        "EndDate": None,
        "IsGlobal": False,
        "IsPublished": True,
    },
]
DROPBOX: list[dict[str, Any]] = [
    {
        "Id": 1,
        "Name": "Homework 1",
        "CustomInstructions": {
            "Text": "Submit your solution as a single PDF.",
            "Html": "<p>Submit your solution as a single PDF.</p>",
        },
        "Availability": None,
        "GroupTypeId": None,
        "IsHidden": False,
    }
]


# --------------------------------------------------------------------------
# Real-format asset generation (module import time): a handful of small
# binaries the /file endpoint serves, picked per ToC entry by that entry's
# own URL extension -- so the bytes served for a given topic are always in
# the format ingest/extract.py will actually try to parse.
#
# `@lru_cache`d, and load-bearing, not just an optimization: PyMuPDF's
# `tobytes()`/python-pptx's `save()` are NOT byte-deterministic across two
# separate calls even with identical inserted text (PDF/OOXML containers
# embed a generation timestamp/id) -- without memoizing per title, two
# fetches of the *same* topic's file (e.g. an E2E run's second /toc
# re-fetching a topic whose D2L LastModifiedDate is null, which is always
# "needed" -- see ingest/diff.py's compute_needed) would hash to two
# different sha256s and look like a genuine content change, defeating both
# the incremental-sync dedupe and the "identical re-run is a pipeline
# no-op" invariant `make e2e` checks.
# --------------------------------------------------------------------------


@lru_cache(maxsize=None)
def make_pdf_bytes(title: str) -> bytes:
    import fitz  # PyMuPDF

    doc = fitz.open()
    try:
        page = doc.new_page()
        page.insert_text((72, 72), title)
        page.insert_text((72, 100), f"Fake D2L fixture content for {title!r}.")
        return doc.tobytes()
    finally:
        doc.close()


@lru_cache(maxsize=None)
def make_pptx_bytes(title: str) -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    body.text_frame.text = f"Fake D2L fixture content for {title!r}."
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@lru_cache(maxsize=None)
def make_vtt_bytes(title: str) -> bytes:
    text = (
        "WEBVTT\n\n"
        "1\n"
        "00:00:00.000 --> 00:00:02.000\n"
        f"{title}\n\n"
        "2\n"
        "00:00:02.000 --> 00:00:04.000\n"
        "Fake D2L fixture caption content.\n"
    )
    return text.encode("utf-8")


@lru_cache(maxsize=None)
def make_html_bytes(title: str) -> bytes:
    html = f"<html><body><h1>{title}</h1><p>Fake D2L fixture HTML content.</p></body></html>"
    return html.encode("utf-8")


def _make_fallback_bytes(title: str) -> bytes:
    return f"Fake D2L fixture content for {title!r}.".encode("utf-8")


_GENERATORS: dict[str, tuple[Callable[[str], bytes], str]] = {
    "pdf": (make_pdf_bytes, "application/pdf"),
    "pptx": (make_pptx_bytes, "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    "vtt": (make_vtt_bytes, "text/vtt"),
    "html": (make_html_bytes, "text/html"),
    "htm": (make_html_bytes, "text/html"),
}
_FALLBACK: tuple[Callable[[str], bytes], str] = (_make_fallback_bytes, "text/plain")


def _extension(value: str | None) -> str:
    if not value:
        return ""
    path = urlparse(value).path or value
    return PurePosixPath(path).suffix.lstrip(".").lower()


def asset_for(title: str, url: str | None) -> tuple[bytes, str]:
    """(content_bytes, content_type) for a topic, picked off `url`'s
    extension (falling back to `title`'s, then to a generic text/plain
    stand-in for anything unrecognized -- never raises)."""
    ext = _extension(url) or _extension(title)
    generator, content_type = _GENERATORS.get(ext, _FALLBACK)
    return generator(title), content_type


# --------------------------------------------------------------------------
# ToC indexing -- topic_id -> {title, url}, walked the same way
# ingest/diff.py's parse_toc walks the real thing (but we only need enough
# to serve a file per topic, not the full entry/module-chain shape).
# --------------------------------------------------------------------------


def _index_topics(toc: dict) -> dict[int, dict[str, Any]]:
    index: dict[int, dict[str, Any]] = {}

    def walk(modules: object) -> None:
        if not isinstance(modules, list):
            return
        for module in modules:
            if not isinstance(module, dict):
                continue
            for topic in module.get("Topics") or []:
                if not isinstance(topic, dict):
                    continue
                topic_id = topic.get("TopicId")
                title = topic.get("Title")
                if isinstance(topic_id, int) and isinstance(title, str):
                    index[topic_id] = {"title": title, "url": topic.get("Url")}
            walk(module.get("Modules"))

    walk(toc.get("Modules"))
    return index


# --------------------------------------------------------------------------
# App factory
# --------------------------------------------------------------------------


def make_fake_d2l(
    fixture_dir: Path,
    *,
    rate_limit_429_every: int | None = None,
    expire_session_after: int | None = None,
) -> FastAPI:
    """Build a fake D2L tenant app serving one course (`ORG_UNIT_ID`).

    `fixture_dir` must contain `toc_sample.json` (reused verbatim as that
    course's ToC, per the Task 13 brief -- existing ingest-API tests already
    pin its exact shape, so this module never modifies it).
    """
    fixture_dir = Path(fixture_dir)
    toc = json.loads((fixture_dir / "toc_sample.json").read_text())
    topic_index = _index_topics(toc)

    app = FastAPI()
    app.state.request_count = 0
    app.state.session_expired = False

    def _fault_response(request: Request) -> Response | None:
        """Shared fault injection for every route below except /_control/*:
        session expiry first (sticky -- once tripped, stays tripped until
        /_control/reset), then periodic rate limiting. Returns a Response to
        short-circuit the route, or None to let it proceed normally."""
        if request.url.path.startswith("/_control"):
            return None

        app.state.request_count += 1

        if expire_session_after is not None and app.state.request_count > expire_session_after:
            app.state.session_expired = True
        if app.state.session_expired:
            return JSONResponse({"detail": "session expired"}, status_code=401)

        if rate_limit_429_every is not None and app.state.request_count % rate_limit_429_every == 0:
            return Response(
                status_code=429,
                headers={
                    "Retry-After": "1",
                    "X-Rate-Limit-Remaining": "0",
                    "X-Rate-Limit-Reset": "1",
                },
            )
        return None

    @app.middleware("http")
    async def _fault_injection(request: Request, call_next):  # type: ignore[no-untyped-def]
        short_circuit = _fault_response(request)
        if short_circuit is not None:
            return short_circuit
        return await call_next(request)

    # ----------------------------------------------------------------------
    # Control endpoints
    # ----------------------------------------------------------------------

    @app.post("/_control/reset")
    def control_reset() -> dict:
        app.state.request_count = 0
        app.state.session_expired = False
        return {"reset": True}

    @app.get("/_control/state")
    def control_state() -> dict:
        return {
            "requestCount": app.state.request_count,
            "sessionExpired": app.state.session_expired,
        }

    # ----------------------------------------------------------------------
    # D2L Valence API surface
    # ----------------------------------------------------------------------

    @app.get("/d2l/api/versions/")
    def versions() -> list[dict]:
        return [
            {"ProductCode": "lp", "LatestVersion": LP_VERSION, "SupportedVersions": [LP_VERSION]},
            {"ProductCode": "le", "LatestVersion": LE_VERSION, "SupportedVersions": [LE_VERSION]},
        ]

    @app.get("/d2l/api/lp/{version}/users/whoami")
    def whoami(version: str) -> dict:
        del version
        return WHOAMI

    @app.get("/d2l/api/lp/{version}/enrollments/myenrollments/")
    def enrollments(version: str, bookmark: str | None = None) -> dict:
        del version
        start = int(bookmark) if bookmark else 0
        page = ENROLLMENTS[start : start + _ENROLLMENT_PAGE_SIZE]
        end = start + len(page)
        has_more = end < len(ENROLLMENTS)
        return {"PagingInfo": {"Bookmark": str(end), "HasMoreItems": has_more}, "Items": page}

    @app.get("/d2l/api/le/{version}/{org_unit_id}/content/toc")
    def content_toc(version: str, org_unit_id: int) -> dict:
        del version
        if org_unit_id != ORG_UNIT_ID:
            raise HTTPException(status_code=404, detail="unknown org unit")
        return toc

    @app.get("/d2l/api/le/{version}/{org_unit_id}/content/topics/{topic_id}/file")
    def content_file(version: str, org_unit_id: int, topic_id: int) -> Response:
        del version, org_unit_id
        entry = topic_index.get(topic_id)
        if entry is None:
            raise HTTPException(status_code=404, detail="unknown topic")
        content, content_type = asset_for(entry["title"], entry["url"])
        return Response(content=content, media_type=content_type)

    @app.get("/d2l/api/le/{version}/{org_unit_id}/news/")
    def news(version: str, org_unit_id: int) -> list[dict]:
        del version, org_unit_id
        return NEWS

    @app.get("/d2l/api/le/{version}/{org_unit_id}/dropbox/folders/")
    def dropbox_folders(version: str, org_unit_id: int) -> list[dict]:
        del version, org_unit_id
        return DROPBOX

    return app
