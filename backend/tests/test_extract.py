"""Tests for text extraction across supported material formats.

Fixtures are generated programmatically (no binary files committed to the
repo): pymupdf builds the PDF, python-pptx the slide deck, python-docx the
Word doc; HTML/VTT/txt are plain strings written to disk.
"""

from __future__ import annotations

import fitz  # PyMuPDF
import pytest
from docx import Document
from pptx import Presentation

from brightspace_agent.ingest.extract import extract_text


# --------------------------------------------------------------------------
# Fixture builders
# --------------------------------------------------------------------------


def _make_pdf(path, pages: list[str]) -> None:
    doc = fitz.open()
    try:
        for text in pages:
            page = doc.new_page()
            page.insert_text((72, 72), text)
        doc.save(path)
    finally:
        doc.close()


def _make_pptx(path) -> None:
    prs = Presentation()
    layout = prs.slide_layouts[1]  # title + content, has a body placeholder

    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Intro to Recursion"
    slide1.placeholders[1].text = "Base case and recursive case"
    slide1.notes_slide.notes_text_frame.text = "Remember to mention stack overflow risk"

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Memoization"
    slide2.placeholders[1].text = "Caching subproblem results"

    prs.save(path)


def _make_docx(path) -> None:
    doc = Document()
    doc.add_paragraph("Course Syllabus")
    doc.add_paragraph("This course covers algorithms and data structures.")
    doc.save(path)


# --------------------------------------------------------------------------
# PDF
# --------------------------------------------------------------------------


def test_extract_pdf_joins_pages_with_form_feed(tmp_path):
    path = tmp_path / "lecture.pdf"
    _make_pdf(path, ["Page one content about graphs", "Page two content about trees"])

    text = extract_text(path, "application/pdf", "document")

    assert text is not None
    assert "graphs" in text
    assert "trees" in text
    assert "\f" in text


def test_extract_pdf_malformed_bytes_returns_none(tmp_path):
    path = tmp_path / "broken.pdf"
    path.write_bytes(b"%PDF-1.4 this is not a real pdf structure at all")

    assert extract_text(path, "application/pdf", "document") is None


# --------------------------------------------------------------------------
# PPTX
# --------------------------------------------------------------------------


def test_extract_pptx_includes_slide_labels_shape_text_and_notes(tmp_path):
    path = tmp_path / "deck.pptx"
    _make_pptx(path)

    text = extract_text(
        path,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "slides",
    )

    assert text is not None
    assert "[Slide 1]" in text
    assert "[Slide 2]" in text
    assert "Intro to Recursion" in text
    assert "Base case and recursive case" in text
    assert "stack overflow risk" in text  # speaker notes
    assert "Memoization" in text


# --------------------------------------------------------------------------
# DOCX
# --------------------------------------------------------------------------


def test_extract_docx_joins_paragraph_text(tmp_path):
    path = tmp_path / "syllabus.docx"
    _make_docx(path)

    text = extract_text(
        path,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "syllabus",
    )

    assert text is not None
    assert "Course Syllabus" in text
    assert "algorithms and data structures" in text


# --------------------------------------------------------------------------
# HTML
# --------------------------------------------------------------------------


def test_extract_html_strips_tags(tmp_path):
    path = tmp_path / "announcement.html"
    path.write_text(
        "<html><body><h1>Homework 3 posted</h1><p>Due <b>Friday</b> at noon.</p></body></html>",
        encoding="utf-8",
    )

    text = extract_text(path, "text/html", "announcement")

    assert text is not None
    assert "Homework 3 posted" in text
    assert "Due Friday at noon." in text
    assert "<h1>" not in text
    assert "<b>" not in text


# --------------------------------------------------------------------------
# VTT / SRT
# --------------------------------------------------------------------------


def test_extract_vtt_strips_header_cue_numbers_and_timestamps_dedupes(tmp_path):
    path = tmp_path / "captions.vtt"
    path.write_text(
        "WEBVTT\n"
        "\n"
        "1\n"
        "00:00:00.000 --> 00:00:02.000\n"
        "Welcome to the lecture on sorting algorithms.\n"
        "\n"
        "2\n"
        "00:00:02.000 --> 00:00:04.000\n"
        "Welcome to the lecture on sorting algorithms.\n"
        "\n"
        "3\n"
        "00:00:04.000 --> 00:00:06.000\n"
        "Today we cover quicksort and mergesort.\n",
        encoding="utf-8",
    )

    text = extract_text(path, "text/vtt", "transcript")

    assert text is not None
    assert "WEBVTT" not in text
    assert "-->" not in text
    assert "quicksort" in text
    # dedup: the identical consecutive cue line should appear only once
    assert text.count("Welcome to the lecture on sorting algorithms.") == 1


def test_extract_srt_strips_cue_numbers_and_timestamps(tmp_path):
    path = tmp_path / "captions.srt"
    path.write_text(
        "1\n"
        "00:00:00,000 --> 00:00:02,000\n"
        "Hello and welcome.\n"
        "\n"
        "2\n"
        "00:00:02,000 --> 00:00:04,000\n"
        "Let's talk about binary search.\n",
        encoding="utf-8",
    )

    text = extract_text(path, "application/x-subrip", "transcript")

    assert text is not None
    assert "-->" not in text
    assert "binary search" in text


# --------------------------------------------------------------------------
# TXT / MD
# --------------------------------------------------------------------------


def test_extract_txt_reads_utf8_and_replaces_bad_bytes(tmp_path):
    path = tmp_path / "notes.txt"
    # Valid UTF-8 content plus one invalid byte sequence spliced in.
    path.write_bytes("Notes on hashing\n".encode("utf-8") + b"\xff\xfe" + "more text".encode("utf-8"))

    text = extract_text(path, "text/plain", "document")

    assert text is not None
    assert "Notes on hashing" in text
    assert "more text" in text  # didn't raise on the bad byte


def test_extract_normalizes_excess_blank_lines(tmp_path):
    path = tmp_path / "notes.md"
    path.write_text("Line one\n\n\n\n\nLine two", encoding="utf-8")

    text = extract_text(path, "text/markdown", "document")

    assert text is not None
    assert "\n\n\n" not in text
    assert "Line one" in text
    assert "Line two" in text


# --------------------------------------------------------------------------
# Unsupported formats
# --------------------------------------------------------------------------


@pytest.mark.parametrize(
    "filename,mime,kind",
    [
        ("lecture.mp4", "video/mp4", "video"),
        ("photo.png", "image/png", "other"),
        ("mystery.bin", "application/x-totally-unknown", "other"),
        ("no_extension", None, "other"),
    ],
)
def test_extract_unsupported_formats_return_none(tmp_path, filename, mime, kind):
    path = tmp_path / filename
    path.write_bytes(b"some opaque binary content")

    assert extract_text(path, mime, kind) is None


def test_extract_missing_file_returns_none(tmp_path):
    path = tmp_path / "does_not_exist.pdf"

    assert extract_text(path, "application/pdf", "document") is None
